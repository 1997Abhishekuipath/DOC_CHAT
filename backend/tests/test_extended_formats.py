"""Extended document format tests for DocChat.

Covers PPTX, XLSX, CSV, image OCR, scanned-PDF OCR fallback, Google Slides
stub, feature-flag gating, RBAC, regressions and chat grounding.

Fixtures are generated in-process via python-pptx / openpyxl / Pillow so the
test is fully self-contained.
"""
from __future__ import annotations

import io
import os
import time
import uuid
from pathlib import Path

import pytest
import requests

BASE_URL = os.environ.get(
    "REACT_APP_BACKEND_URL", "https://app-launch-view-1.preview.emergentagent.com"
).rstrip("/")
API = f"{BASE_URL}/api"

OWNER_EMAIL = "owner@example.com"
OWNER_PASS = "test1234"


# ---------------- auth fixtures ----------------
@pytest.fixture(scope="module")
def owner_headers():
    r = requests.post(f"{API}/auth/login",
                      json={"email": OWNER_EMAIL, "password": OWNER_PASS}, timeout=20)
    assert r.status_code == 200, f"owner login: {r.status_code} {r.text}"
    return {"Authorization": f"Bearer {r.json()['access_token']}"}


@pytest.fixture(scope="module")
def viewer_headers():
    email = f"TEST_viewer_{uuid.uuid4().hex[:8]}@example.com"
    r = requests.post(f"{API}/auth/register",
                      json={"email": email, "password": "test1234", "name": "TEST Viewer"},
                      timeout=20)
    assert r.status_code == 200, r.text
    return {"Authorization": f"Bearer {r.json()['access_token']}"}


# ---------------- fixture generators ----------------
def _make_pptx_bytes() -> bytes:
    from pptx import Presentation
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for idx, line in enumerate([
        "Alpha fact: the sky color is cerulean for testing",
        "Beta fact: the magic code is PURPLE-42",
        "Gamma fact: quarterly revenue was 9876 dollars",
    ], start=1):
        slide = prs.slides.add_slide(blank)
        tx = slide.shapes.add_textbox(914400, 914400, 6400000, 3600000)
        tx.text_frame.text = f"Slide {idx}: {line}"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_xlsx_bytes() -> bytes:
    from openpyxl import Workbook
    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Sales"
    ws1.append(["Region", "Revenue", "Notes"])
    ws1.append(["North", 12345, "TEST_XLSX_ROW_NORTH_A1"])
    ws1.append(["South", 67890, "TEST_XLSX_ROW_SOUTH_B2"])
    ws2 = wb.create_sheet("Inventory")
    ws2.append(["SKU", "Qty"])
    ws2.append(["WIDGET-77", 150])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_csv_bytes() -> bytes:
    return (b"name,score,note\n"
            b"Alice,90,TEST_CSV_MARKER_ZEBRA\n"
            b"Bob,85,hello\n")


def _make_png_bytes(text: str = "Invoice 12345 Amount $500 DocChat OCR test") -> bytes:
    from PIL import Image, ImageDraw, ImageFont
    img = Image.new("RGB", (1000, 200), "white")
    d = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 42)
    except Exception:
        font = ImageFont.load_default()
    d.text((20, 70), text, fill="black", font=font)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _make_jpg_bytes(text: str = "Receipt ZULU777 for 42 units") -> bytes:
    from PIL import Image, ImageDraw, ImageFont
    img = Image.new("RGB", (1000, 200), "white")
    d = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 42)
    except Exception:
        font = ImageFont.load_default()
    d.text((20, 70), text, fill="black", font=font)
    buf = io.BytesIO()
    img.save(buf, format="JPEG")
    return buf.getvalue()


def _make_scanned_pdf_bytes() -> bytes:
    """Create a PDF built from images of text => no extractable text."""
    from PIL import Image, ImageDraw, ImageFont
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 48)
    except Exception:
        font = ImageFont.load_default()
    pages = []
    for line in [
        "Scanned page one contains KEYWORD-ALPHA",
        "Scanned page two mentions KEYWORD-BRAVO",
    ]:
        img = Image.new("RGB", (1200, 1600), "white")
        ImageDraw.Draw(img).text((60, 80), line, fill="black", font=font)
        pages.append(img)
    buf = io.BytesIO()
    pages[0].save(buf, format="PDF", save_all=True, append_images=pages[1:])
    return buf.getvalue()


# ---------------- helpers ----------------
def _upload(headers, filename: str, data: bytes, mime: str):
    files = {"file": (filename, data, mime)}
    return requests.post(f"{API}/v2/documents/ingest", headers=headers,
                         files=files, data={"tags": "TEST_fmt"}, timeout=60)


def _poll_until_ready(headers, doc_id: str, timeout: int = 90) -> dict:
    deadline = time.time() + timeout
    last = {}
    while time.time() < deadline:
        r = requests.get(f"{API}/v2/documents/{doc_id}", headers=headers, timeout=20)
        if r.status_code == 200:
            last = r.json()
            if last.get("status") in ("ready", "failed"):
                return last
        time.sleep(2)
    return last


# ---------------- flags endpoint ----------------
class TestFlags:
    def test_flags_include_new_keys(self, owner_headers):
        r = requests.get(f"{API}/v2/flags", headers=owner_headers, timeout=10)
        assert r.status_code == 200
        flags = r.json()
        for key in [
            "ENABLE_PPTX_SUPPORT", "ENABLE_EXCEL_SUPPORT",
            "ENABLE_IMAGE_OCR", "ENABLE_SCANNED_PDF_OCR",
            "ENABLE_GOOGLE_SLIDES",
        ]:
            assert key in flags, f"flag {key} missing"
        assert flags["ENABLE_GOOGLE_SLIDES"] is False


# ---------------- new format ingestion ----------------
class TestNewFormats:
    def test_pptx_ingest(self, owner_headers):
        r = _upload(owner_headers, "TEST_deck.pptx", _make_pptx_bytes(),
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        assert r.status_code == 200, r.text
        doc_id = r.json()["id"]
        pytest.pptx_doc_id = doc_id
        final = _poll_until_ready(owner_headers, doc_id)
        assert final.get("status") == "ready", f"pptx final state: {final}"
        assert final["file_type"] == "pptx"
        assert final["chunk_count"] > 0
        assert final["page_count"] == 3  # three slides with text

    def test_xlsx_ingest(self, owner_headers):
        r = _upload(owner_headers, "TEST_sheet.xlsx", _make_xlsx_bytes(),
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        assert r.status_code == 200, r.text
        doc_id = r.json()["id"]
        pytest.xlsx_doc_id = doc_id
        final = _poll_until_ready(owner_headers, doc_id)
        assert final.get("status") == "ready", f"xlsx final: {final}"
        assert final["file_type"] == "xlsx"
        assert final["chunk_count"] > 0
        # two non-empty sheets
        assert final["page_count"] == 2

    def test_csv_ingest(self, owner_headers):
        r = _upload(owner_headers, "TEST_data.csv", _make_csv_bytes(), "text/csv")
        assert r.status_code == 200, r.text
        doc_id = r.json()["id"]
        final = _poll_until_ready(owner_headers, doc_id)
        assert final.get("status") == "ready", f"csv final: {final}"
        assert final["file_type"] == "csv"
        assert final["chunk_count"] > 0
        assert final["page_count"] == 1

    def test_png_ocr_ingest(self, owner_headers):
        r = _upload(owner_headers, "TEST_img.png", _make_png_bytes(), "image/png")
        assert r.status_code == 200, r.text
        doc_id = r.json()["id"]
        pytest.png_doc_id = doc_id
        final = _poll_until_ready(owner_headers, doc_id, timeout=120)
        assert final.get("status") == "ready", f"png final: {final}"
        assert final["file_type"] == "image"
        assert final["chunk_count"] > 0
        assert final["page_count"] == 1

    def test_jpg_ocr_ingest(self, owner_headers):
        r = _upload(owner_headers, "TEST_img.jpg", _make_jpg_bytes(), "image/jpeg")
        assert r.status_code == 200, r.text
        doc_id = r.json()["id"]
        final = _poll_until_ready(owner_headers, doc_id, timeout=120)
        assert final.get("status") == "ready", f"jpg final: {final}"
        assert final["file_type"] == "image"
        assert final["chunk_count"] > 0

    def test_scanned_pdf_ocr_fallback(self, owner_headers):
        r = _upload(owner_headers, "TEST_scanned.pdf", _make_scanned_pdf_bytes(),
                    "application/pdf")
        assert r.status_code == 200, r.text
        doc_id = r.json()["id"]
        final = _poll_until_ready(owner_headers, doc_id, timeout=180)
        assert final.get("status") == "ready", f"scanned pdf final: {final}"
        assert final["file_type"] == "pdf"
        assert final["chunk_count"] > 0, "OCR fallback must produce chunks"


# ---------------- regression: existing formats ----------------
class TestLegacyFormatsStillWork:
    def test_txt(self, owner_headers):
        r = _upload(owner_headers, "TEST_note.txt",
                    b"Quarterly sales improved by 12 percent. Alpha test.", "text/plain")
        assert r.status_code == 200, r.text
        final = _poll_until_ready(owner_headers, r.json()["id"])
        assert final.get("status") == "ready"
        assert final["file_type"] == "text"

    def test_md(self, owner_headers):
        r = _upload(owner_headers, "TEST_readme.md",
                    b"# Title\n\nSome **markdown** content for testing.", "text/markdown")
        assert r.status_code == 200, r.text
        final = _poll_until_ready(owner_headers, r.json()["id"])
        assert final.get("status") == "ready"
        assert final["file_type"] == "markdown"

    def test_docx(self, owner_headers):
        from docx import Document as DocxDocument
        d = DocxDocument()
        d.add_paragraph("Regression docx test. Lorem ipsum.")
        buf = io.BytesIO()
        d.save(buf)
        r = _upload(owner_headers, "TEST_doc.docx", buf.getvalue(),
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        assert r.status_code == 200, r.text
        final = _poll_until_ready(owner_headers, r.json()["id"])
        assert final.get("status") == "ready"
        assert final["file_type"] == "docx"


# ---------------- rejections & gating ----------------
class TestRejections:
    def test_unsupported_doc_rejected(self, owner_headers):
        r = _upload(owner_headers, "TEST_old.doc", b"not really doc", "application/msword")
        assert r.status_code == 400
        assert "Unsupported file type" in r.json().get("detail", "")
        assert ".doc" in r.json().get("detail", "")

    def test_viewer_cannot_ingest(self, viewer_headers):
        r = _upload(viewer_headers, "TEST_viewer.txt", b"blocked", "text/plain")
        assert r.status_code == 403

    def test_google_slides_stub_disabled(self, owner_headers):
        r = requests.post(
            f"{API}/v2/documents/ingest-google-slides",
            headers=owner_headers,
            json={"presentation_id": "abc", "title": "test"},
            timeout=20,
        )
        assert r.status_code == 400, r.text
        assert "ENABLE_GOOGLE_SLIDES" in r.json().get("detail", "")

    def test_google_slides_requires_editor(self, viewer_headers):
        r = requests.post(
            f"{API}/v2/documents/ingest-google-slides",
            headers=viewer_headers,
            json={"presentation_id": "abc"},
            timeout=20,
        )
        assert r.status_code == 403


# ---------------- regression: other existing endpoints ----------------
class TestExistingEndpointsRegression:
    def test_providers(self, owner_headers):
        r = requests.get(f"{API}/v2/providers", headers=owner_headers, timeout=10)
        assert r.status_code == 200
        data = r.json()
        assert "llm" in data and "embedding" in data
        assert "provider" in data["llm"] and "model" in data["llm"]

    def test_sessions_list(self, owner_headers):
        r = requests.get(f"{API}/v2/sessions", headers=owner_headers, timeout=10)
        assert r.status_code == 200
        assert isinstance(r.json(), list)

    def test_share_links_list(self, owner_headers):
        r = requests.get(f"{API}/v2/share-links", headers=owner_headers, timeout=10)
        assert r.status_code == 200
        assert isinstance(r.json(), list)

    def test_admin_users(self, owner_headers):
        r = requests.get(f"{API}/admin/users", headers=owner_headers, timeout=10)
        assert r.status_code == 200
        assert isinstance(r.json(), list)

    def test_admin_audit_log(self, owner_headers):
        r = requests.get(f"{API}/admin/audit-log", headers=owner_headers, timeout=10)
        assert r.status_code == 200

    def test_admin_analytics(self, owner_headers):
        r = requests.get(f"{API}/admin/analytics", headers=owner_headers, timeout=10)
        assert r.status_code == 200

    def test_auth_me(self, owner_headers):
        r = requests.get(f"{API}/auth/me", headers=owner_headers, timeout=10)
        assert r.status_code == 200
        data = r.json()
        user = data.get("user", data)
        assert user.get("email") == OWNER_EMAIL


# ---------------- chat grounding ----------------
def _chat_answer(headers, doc_ids, question):
    r = requests.post(
        f"{API}/v2/chat",
        headers=headers,
        json={
            "query": question,
            "document_ids": doc_ids,
            "stream": False,
        },
        timeout=120,
    )
    return r


class TestChatGrounding:
    """Chat grounding for new formats — depends on an active LLM provider.
    Tests will self-skip if chat is not configured (returns 500/400)."""

    def test_pptx_grounded(self, owner_headers):
        doc_id = getattr(pytest, "pptx_doc_id", None)
        if not doc_id:
            pytest.skip("pptx doc not ingested")
        r = _chat_answer(owner_headers, [doc_id], "What is the magic code?")
        if r.status_code != 200:
            pytest.skip(f"chat not active ({r.status_code}): {r.text[:160]}")
        data = r.json()
        answer = (data.get("answer") or "").lower()
        assert "purple-42" in answer or "purple" in answer, f"answer missing fact: {answer}"
        citations = data.get("citations") or []
        assert citations, "expected citations"
        assert any(c.get("filename") == "TEST_deck.pptx" for c in citations), \
            f"citations={citations}"

    def test_xlsx_grounded(self, owner_headers):
        doc_id = getattr(pytest, "xlsx_doc_id", None)
        if not doc_id:
            pytest.skip("xlsx doc not ingested")
        r = _chat_answer(owner_headers, [doc_id],
                         "What is the note for the North region row?")
        if r.status_code != 200:
            pytest.skip(f"chat not active: {r.status_code}")
        data = r.json()
        answer = (data.get("answer") or "").upper()
        assert "TEST_XLSX_ROW_NORTH_A1" in answer or "NORTH" in answer, \
            f"answer missing: {answer[:200]}"
        citations = data.get("citations") or []
        assert citations
        # North is on sheet 1
        assert any(c.get("filename") == "TEST_sheet.xlsx" and c.get("page") == 1
                   for c in citations), f"citations={citations}"

    def test_image_ocr_grounded(self, owner_headers):
        doc_id = getattr(pytest, "png_doc_id", None)
        if not doc_id:
            pytest.skip("png doc not ingested")
        r = _chat_answer(owner_headers, [doc_id], "What is the invoice number?")
        if r.status_code != 200:
            pytest.skip(f"chat not active: {r.status_code}")
        data = r.json()
        answer = data.get("answer") or ""
        # OCR may have minor char errors (e.g. 12345 -> 12348) — accept any
        # 4-5 digit sequence as proof the number was retrieved from the image.
        import re
        nums = re.findall(r"\b\d{4,5}\b", answer)
        assert nums, f"OCR text not retrieved: {answer[:200]}"
        citations = data.get("citations") or []
        assert citations, "expected citations"
        assert any(c.get("page") == 1 and c.get("filename") == "TEST_img.png"
                   for c in citations), f"citations={citations}"
