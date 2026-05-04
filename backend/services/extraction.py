"""Document text extraction — multi-modal enhanced pipeline.

Supported formats (feature-flag gated for newer types):
  - PDF        (always — native text + optional table/image-OCR layers)
  - DOCX       (always — paragraphs + tables)
  - TXT / MD   (always)
  - PPTX       (ENABLE_PPTX_SUPPORT — shape text + optional image OCR)
  - XLSX / CSV (ENABLE_EXCEL_SUPPORT)
  - Images     (ENABLE_IMAGE_OCR)  — .png / .jpg / .jpeg

Advanced extraction flags (all default ON; set to false to downgrade):
  ENABLE_ADVANCED_OCR        — multi-pass OCR with preprocessing
  ENABLE_TABLE_EXTRACTION    — structured table extraction from PDF/PPTX
  ENABLE_IMAGE_IN_PDF_OCR    — OCR on embedded images inside PDF pages
  ENABLE_PPTX_IMAGE_OCR      — OCR on picture shapes inside PPTX slides
  ENABLE_SCANNED_PDF_OCR     — full-page OCR for scanned/empty PDF pages

Each returned unit is (unit_number, merged_text) where unit_number is:
  - page number  for PDF
  - slide number for PPTX
  - sheet index  for XLSX (1-based)
  - always 1     for DOCX / TXT / MD / CSV / image

Merged text uses lightweight section markers:
  [TABLE]       — structured table rows follow
  [IMAGE_TEXT]  — OCR text from an embedded image follows

These markers make individual content types queryable inside RAG without
changing any downstream API.
"""
from __future__ import annotations

import csv as _csv
import io
import logging
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from core.config import is_enabled

logger = logging.getLogger("docchat.extraction")

# ── Extension sets ────────────────────────────────────────────────────────────
_CORE_EXTS   = {".pdf", ".docx", ".txt", ".md"}
_PPTX_EXTS   = {".pptx"}
_EXCEL_EXTS  = {".xlsx", ".csv"}
_IMAGE_EXTS  = {".png", ".jpg", ".jpeg"}
_GSLIDES_EXTS: set[str] = set()

ALL_KNOWN_EXTENSIONS = _CORE_EXTS | _PPTX_EXTS | _EXCEL_EXTS | _IMAGE_EXTS


class UnsupportedFormatError(ValueError):
    """Raised when a file type is not enabled or not supported."""


def current_supported_extensions() -> set[str]:
    exts = set(_CORE_EXTS)
    if is_enabled("ENABLE_PPTX_SUPPORT"):
        exts |= _PPTX_EXTS
    if is_enabled("ENABLE_EXCEL_SUPPORT"):
        exts |= _EXCEL_EXTS
    if is_enabled("ENABLE_IMAGE_OCR"):
        exts |= _IMAGE_EXTS
    return exts


SUPPORTED_EXTENSIONS = current_supported_extensions()


def file_type_label(filename: str) -> str:
    ext = Path(filename).suffix.lower()
    return {
        ".pdf": "pdf", ".docx": "docx", ".pptx": "pptx",
        ".xlsx": "xlsx", ".csv": "csv",
        ".png": "image", ".jpg": "image", ".jpeg": "image",
        ".md": "markdown", ".txt": "text",
    }.get(ext, ext.lstrip(".") or "unknown")


# ── Public dispatcher ─────────────────────────────────────────────────────────
def extract_text(file_path: Path, filename: str) -> List[Tuple[int, str]]:
    """Dispatch to the right extractor. Raises UnsupportedFormatError on unknown/disabled types."""
    ext = Path(filename).suffix.lower()

    if ext == ".pdf":
        return _extract_pdf(file_path)
    if ext == ".docx":
        return _extract_docx(file_path)
    if ext in (".txt", ".md"):
        return _extract_plaintext(file_path)
    if ext == ".pptx":
        if not is_enabled("ENABLE_PPTX_SUPPORT"):
            raise UnsupportedFormatError("PPTX support is disabled. Enable ENABLE_PPTX_SUPPORT.")
        return _extract_pptx(file_path)
    if ext == ".xlsx":
        if not is_enabled("ENABLE_EXCEL_SUPPORT"):
            raise UnsupportedFormatError("Excel support is disabled. Enable ENABLE_EXCEL_SUPPORT.")
        return _extract_xlsx(file_path)
    if ext == ".csv":
        if not is_enabled("ENABLE_EXCEL_SUPPORT"):
            raise UnsupportedFormatError("CSV support is disabled. Enable ENABLE_EXCEL_SUPPORT.")
        return _extract_csv(file_path)
    if ext in _IMAGE_EXTS:
        if not is_enabled("ENABLE_IMAGE_OCR"):
            raise UnsupportedFormatError("Image OCR is disabled. Enable ENABLE_IMAGE_OCR.")
        return _extract_image_ocr(file_path)

    raise UnsupportedFormatError(f"Unsupported file type: {ext}")


# ═══════════════════════════════════════════════════════════════════════════════
# OCR helpers — preprocessing + multi-pass
# ═══════════════════════════════════════════════════════════════════════════════

def _preprocess_for_ocr(pil_image):
    """Apply grayscale → upscale-if-tiny → Otsu threshold → return preprocessed PIL image."""
    try:
        import numpy as np
        from PIL import Image, ImageFilter

        img = pil_image.convert("L")

        # Upscale very small images to help Tesseract
        w, h = img.size
        if min(w, h) < 100 or max(w, h) < 300:
            scale = max(2, int(300 / max(min(w, h), 1)))
            scale = min(scale, 4)  # cap at 4×
            img = img.resize((w * scale, h * scale), Image.LANCZOS)

        # Mild sharpening before thresholding
        img = img.filter(ImageFilter.SHARPEN)

        # Otsu-style global threshold via numpy
        arr = np.array(img, dtype=np.float32)
        thresh = arr.mean()
        binary = ((arr >= thresh) * 255).astype(np.uint8)

        from PIL import Image as PILImage
        return PILImage.fromarray(binary)
    except Exception:
        # Graceful degradation — return grayscale at minimum
        try:
            return pil_image.convert("L")
        except Exception:
            return pil_image


def _ocr_image(pil_image, source_hint: str = "") -> str:
    """Run OCR on a PIL image.

    If ENABLE_ADVANCED_OCR is on: multi-pass with preprocessing.
    Otherwise: single-pass grayscale (original behaviour).
    Returns stripped text or empty string.
    """
    try:
        import pytesseract
    except ImportError:
        return ""

    if not is_enabled("ENABLE_ADVANCED_OCR"):
        # Original simple path
        try:
            gray = pil_image.convert("L")
            return (pytesseract.image_to_string(gray) or "").strip()
        except Exception:
            return ""

    # ── Multi-pass advanced path ──────────────────────────────────────────────
    results: List[str] = []

    # Pass 1 — preprocessed (Otsu + upscaled)
    try:
        preprocessed = _preprocess_for_ocr(pil_image)
        t = (pytesseract.image_to_string(preprocessed, config="--psm 6") or "").strip()
        if t:
            results.append(t)
    except Exception as e:
        logger.debug("OCR pass 1 failed (%s): %s", source_hint, e)

    # Pass 2 — different page-segmentation mode (auto with OSD)
    if not results:
        try:
            preprocessed = _preprocess_for_ocr(pil_image)
            t = (pytesseract.image_to_string(preprocessed, config="--psm 3") or "").strip()
            if t:
                results.append(t)
        except Exception as e:
            logger.debug("OCR pass 2 failed (%s): %s", source_hint, e)

    # Pass 3 — plain grayscale, original size (fallback)
    if not results:
        try:
            gray = pil_image.convert("L")
            t = (pytesseract.image_to_string(gray, config="--psm 6") or "").strip()
            if t:
                results.append(t)
        except Exception as e:
            logger.debug("OCR pass 3 failed (%s): %s", source_hint, e)

    # Return the longest successful result
    if results:
        return max(results, key=len)
    return ""


# ═══════════════════════════════════════════════════════════════════════════════
# Table helpers
# ═══════════════════════════════════════════════════════════════════════════════

def _format_table_rows(rows: List[List[Optional[str]]]) -> str:
    """Format a list-of-rows (each row = list of cell strings) into clean pipe-delimited text."""
    lines: List[str] = []
    for row in rows:
        cells = [str(c).strip() if c is not None else "" for c in row]
        if any(cells):
            lines.append(" | ".join(cells))
    return "\n".join(lines)


def _extract_pdf_tables(file_path: Path) -> Dict[int, str]:
    """Extract tables per page using pdfplumber. Returns {page_num: table_text}."""
    if not is_enabled("ENABLE_TABLE_EXTRACTION"):
        return {}
    try:
        import pdfplumber
    except ImportError:
        logger.warning("pdfplumber not installed — table extraction skipped")
        return {}

    page_tables: Dict[int, List[str]] = {}
    try:
        with pdfplumber.open(str(file_path)) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                try:
                    tables = page.extract_tables() or []
                except Exception:
                    continue
                for table in tables:
                    if not table:
                        continue
                    formatted = _format_table_rows(table)
                    if formatted.strip():
                        page_tables.setdefault(i, []).append(formatted)
    except Exception as e:
        logger.warning("pdfplumber table extraction failed: %s", e)
        return {}

    return {
        pg: "[TABLE]\n" + "\n\n".join(tbs)
        for pg, tbs in page_tables.items()
        if tbs
    }


def _extract_pdf_embedded_image_ocr(file_path: Path) -> Dict[int, str]:
    """Extract embedded images from PDF pages via PyMuPDF and OCR them.
    Returns {page_num: image_ocr_text}.
    """
    if not is_enabled("ENABLE_IMAGE_IN_PDF_OCR"):
        return {}
    try:
        import fitz  # PyMuPDF
        from PIL import Image
    except ImportError:
        logger.warning("PyMuPDF or Pillow not installed — PDF embedded image OCR skipped")
        return {}

    results: Dict[int, List[str]] = {}
    try:
        doc = fitz.open(str(file_path))
        for page_idx in range(len(doc)):
            page_num = page_idx + 1
            try:
                image_list = doc[page_idx].get_images(full=True)
            except Exception:
                continue

            for img_idx, img_info in enumerate(image_list):
                xref = img_info[0]
                try:
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image.get("image", b"")
                    if not image_bytes:
                        continue
                    pil_img = Image.open(io.BytesIO(image_bytes))
                    hint = f"pdf_p{page_num}_img{img_idx}"
                    ocr_text = _ocr_image(pil_img, source_hint=hint)
                    if ocr_text:
                        results.setdefault(page_num, []).append(ocr_text)
                except Exception as e:
                    logger.debug("PDF image OCR failed (p%d img%d): %s", page_num, img_idx, e)
        doc.close()
    except Exception as e:
        logger.warning("PDF embedded image OCR failed: %s", e)

    return {
        pg: "[IMAGE_TEXT]\n" + "\n\n".join(txts)
        for pg, txts in results.items()
        if txts
    }


# ═══════════════════════════════════════════════════════════════════════════════
# PDF extraction — unified
# ═══════════════════════════════════════════════════════════════════════════════

def _extract_pdf(file_path: Path) -> List[Tuple[int, str]]:
    """Full PDF extraction: native text + tables + embedded image OCR + scanned fallback."""
    from pypdf import PdfReader

    # ── Pre-gather enhanced data (done once per file, not per page) ──────────
    tables_by_page: Dict[int, str] = _extract_pdf_tables(file_path)
    img_ocr_by_page: Dict[int, str] = _extract_pdf_embedded_image_ocr(file_path)

    try:
        reader = PdfReader(str(file_path))
    except Exception as e:
        logger.error("PdfReader failed on %s: %s", file_path.name, e)
        return []

    result: List[Tuple[int, str]] = []
    empty_pages: List[int] = []

    for i, page in enumerate(reader.pages, start=1):
        try:
            native_text = (page.extract_text() or "").strip()
        except Exception:
            native_text = ""

        parts: List[str] = []
        if native_text:
            parts.append(native_text)

        # Append structured table text for this page
        if i in tables_by_page:
            parts.append(tables_by_page[i])

        # Append image OCR for this page
        if i in img_ocr_by_page:
            parts.append(img_ocr_by_page[i])

        merged = "\n\n".join(p for p in parts if p.strip())
        if merged.strip():
            result.append((i, merged))
        else:
            empty_pages.append(i)

    # Scanned PDF fallback: full-page render + OCR for pages with zero text
    if empty_pages and is_enabled("ENABLE_SCANNED_PDF_OCR"):
        ocr_pages = _ocr_pdf_pages(file_path, empty_pages)
        result.extend(ocr_pages)
        result.sort(key=lambda x: x[0])

    return result


def _ocr_pdf_pages(file_path: Path, page_numbers: List[int]) -> List[Tuple[int, str]]:
    """Render selected PDF pages to images and run Tesseract OCR (scanned PDF fallback)."""
    try:
        from pdf2image import convert_from_path
    except ImportError:
        logger.warning("pdf2image not installed — scanned PDF OCR skipped")
        return []

    first = min(page_numbers)
    last = max(page_numbers)
    try:
        images = convert_from_path(str(file_path), dpi=200, first_page=first, last_page=last)
    except Exception as e:
        logger.warning("pdf2image convert failed: %s", e)
        return []

    wanted = set(page_numbers)
    out: List[Tuple[int, str]] = []
    for offset, img in enumerate(images):
        page_num = first + offset
        if page_num not in wanted:
            continue
        text = _ocr_image(img, source_hint=f"scanned_p{page_num}")
        if text:
            out.append((page_num, text))
    return out


# ═══════════════════════════════════════════════════════════════════════════════
# DOCX — paragraphs + structured tables
# ═══════════════════════════════════════════════════════════════════════════════

def _extract_docx(file_path: Path) -> List[Tuple[int, str]]:
    from docx import Document as DocxDocument

    doc = DocxDocument(str(file_path))
    parts: List[str] = []

    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)

    if is_enabled("ENABLE_TABLE_EXTRACTION"):
        for table in doc.tables:
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            formatted = _format_table_rows(rows)
            if formatted.strip():
                parts.append("[TABLE]\n" + formatted)
    else:
        # Original plain-text table fallback
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                parts.append("\t".join(cells))

    text = "\n".join(parts)
    return [(1, text)] if text.strip() else []


# ═══════════════════════════════════════════════════════════════════════════════
# TXT / MD
# ═══════════════════════════════════════════════════════════════════════════════

def _extract_plaintext(file_path: Path) -> List[Tuple[int, str]]:
    text = file_path.read_text(encoding="utf-8", errors="ignore")
    return [(1, text)] if text.strip() else []

# Backwards compat alias
_extract_text = _extract_plaintext


# ═══════════════════════════════════════════════════════════════════════════════
# PPTX — shape text + table shapes + optional image OCR
# ═══════════════════════════════════════════════════════════════════════════════

def _ocr_pptx_shape_image(shape) -> str:
    """Extract and OCR image data from a PPTX picture shape."""
    if not is_enabled("ENABLE_PPTX_IMAGE_OCR"):
        return ""
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER  # noqa: F401
        from PIL import Image
        image_blob = shape.image.blob
        pil_img = Image.open(io.BytesIO(image_blob))
        return _ocr_image(pil_img, source_hint="pptx_img")
    except Exception as e:
        logger.debug("PPTX shape image OCR failed: %s", e)
        return ""


def _extract_pptx(file_path: Path) -> List[Tuple[int, str]]:
    """PPTX extraction: text frames + tables + optional embedded image OCR."""
    from pptx import Presentation

    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        _has_mso = True
    except ImportError:
        _has_mso = False

    prs = Presentation(str(file_path))
    out: List[Tuple[int, str]] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []

        for shape in slide.shapes:
            # ── Text frames ──
            if shape.has_text_frame:
                lines = []
                for para in shape.text_frame.paragraphs:
                    line = "".join(r.text for r in para.runs).strip()
                    if line:
                        lines.append(line)
                if lines:
                    parts.append("\n".join(lines))

            # ── Table shapes ──
            if is_enabled("ENABLE_TABLE_EXTRACTION") and shape.has_table:
                rows = []
                for row in shape.table.rows:
                    cells = [(cell.text or "").strip() for cell in row.cells]
                    rows.append(cells)
                formatted = _format_table_rows(rows)
                if formatted.strip():
                    parts.append("[TABLE]\n" + formatted)

            # ── Picture shapes → OCR ──
            if is_enabled("ENABLE_PPTX_IMAGE_OCR"):
                try:
                    is_picture = (
                        _has_mso and shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                    )
                    if is_picture:
                        img_text = _ocr_pptx_shape_image(shape)
                        if img_text:
                            parts.append("[IMAGE_TEXT]\n" + img_text)
                except Exception as e:
                    logger.debug("PPTX picture shape check failed (slide %d): %s", slide_idx, e)

        # Speaker notes
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    parts.append(f"[Notes] {notes}")
        except Exception:
            pass

        text = "\n\n".join(p for p in parts if p.strip())
        if text:
            out.append((slide_idx, text))

    return out


# ═══════════════════════════════════════════════════════════════════════════════
# XLSX — structured rows; one unit per sheet
# ═══════════════════════════════════════════════════════════════════════════════

def _extract_xlsx(file_path: Path) -> List[Tuple[int, str]]:
    from openpyxl import load_workbook

    wb = load_workbook(filename=str(file_path), read_only=True, data_only=True)
    out: List[Tuple[int, str]] = []

    for idx, sheet_name in enumerate(wb.sheetnames, start=1):
        ws = wb[sheet_name]
        rows: List[List[str]] = []
        for row in ws.iter_rows(values_only=True):
            cells = ["" if c is None else str(c).strip() for c in row]
            if any(cells):
                rows.append(cells)

        if not rows:
            continue

        if is_enabled("ENABLE_TABLE_EXTRACTION"):
            header = f"Sheet: {sheet_name}"
            formatted = _format_table_rows(rows)
            text = f"{header}\n[TABLE]\n{formatted}"
        else:
            lines = [f"Sheet: {sheet_name}"] + ["\t".join(r) for r in rows]
            text = "\n".join(lines)

        if text.strip():
            out.append((idx, text))

    return out


# ═══════════════════════════════════════════════════════════════════════════════
# CSV
# ═══════════════════════════════════════════════════════════════════════════════

def _extract_csv(file_path: Path) -> List[Tuple[int, str]]:
    rows: List[List[str]] = []
    with file_path.open("r", encoding="utf-8", errors="ignore", newline="") as f:
        reader = _csv.reader(f)
        for row in reader:
            if any((c or "").strip() for c in row):
                rows.append(row)

    if not rows:
        return []

    if is_enabled("ENABLE_TABLE_EXTRACTION"):
        formatted = _format_table_rows(rows)
        text = "[TABLE]\n" + formatted
    else:
        text = "\n".join("\t".join(r) for r in rows)

    return [(1, text)] if text.strip() else []


# ═══════════════════════════════════════════════════════════════════════════════
# Image OCR (standalone image files: PNG / JPG / JPEG)
# ═══════════════════════════════════════════════════════════════════════════════

def _extract_image_ocr(file_path: Path) -> List[Tuple[int, str]]:
    """OCR a standalone image file.

    With ENABLE_ADVANCED_OCR: multi-pass preprocessing pipeline.
    Always falls back to rich metadata if OCR yields nothing.
    """
    from PIL import Image

    try:
        pil_img = Image.open(str(file_path))
    except Exception as e:
        logger.warning("Cannot open image %s: %s", file_path.name, e)
        return [(1, f"[Image file: {file_path.stem}] Could not open image.")]

    ocr_text = _ocr_image(pil_img, source_hint=file_path.stem)

    if ocr_text:
        return [(1, ocr_text)]

    # Rich metadata fallback — always indexable, never empty
    try:
        w, h = pil_img.size
        fmt = pil_img.format or "unknown"
        mode = pil_img.mode
        fallback = (
            f"[Image file: {file_path.stem}] "
            f"Format: {fmt}, Size: {w}×{h}px, Mode: {mode}. "
            f"No machine-readable text was detected by OCR. "
            f"The image may be a photograph, diagram, or decorative graphic."
        )
    except Exception:
        fallback = f"[Image file: {file_path.stem}] Image document with no detectable text."

    return [(1, fallback)]


# ═══════════════════════════════════════════════════════════════════════════════
# Google Slides — STUB
# ═══════════════════════════════════════════════════════════════════════════════

def extract_google_slides(presentation_id: str) -> List[Tuple[int, str]]:
    """Stub — requires Google API credentials. Not yet implemented."""
    if not is_enabled("ENABLE_GOOGLE_SLIDES"):
        raise UnsupportedFormatError(
            "Google Slides ingestion is disabled. Enable ENABLE_GOOGLE_SLIDES and configure credentials."
        )
    raise NotImplementedError(
        "Google Slides ingestion is not yet implemented. "
        "Flag is on but credentials are not configured."
    )
